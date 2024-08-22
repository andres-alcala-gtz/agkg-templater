import re
import docx
import pptx
import shutil
import pathlib
import openpyxl

import utilities


def update_xls(path: pathlib.Path, data: dict[str, str]) -> None:

    def _update(container) -> None:
        if isinstance(container.value, str) and re.search("«[\w]+»", container.value) and container.data_type != "f":
            text = container.value
            for key, value in data.items():
                text = text.replace(key, value)
            container.value = text

    workbook = openpyxl.load_workbook(str(path))
    for sheetname, (rows, cols) in utilities.worksheets_dimensions(str(path)).items():
        for row in range(1, rows + 2):
            for col in range(1, cols + 1):
                cell = workbook[sheetname].cell(row, col)
                _update(cell)
    workbook.save(str(path))


def update_doc(path: pathlib.Path, data: dict[str, str]) -> None:

    def _update(container) -> None:
        if isinstance(container.text, str) and re.search("«[\w]+»", container.text):
            concatenate = False
            for index in reversed(range(len(container.runs))):
                if concatenate:
                    container.runs[index].text += container.runs[index + 1].text
                    container.runs[index + 1].text = ""
                if "»" in container.runs[index].text:
                    concatenate = True
                if "«" in container.runs[index].text:
                    concatenate = False
                    container.runs[index].text = data[container.runs[index].text]

    document = docx.Document(str(path))
    paragraphs = document.paragraphs
    tables = document.tables
    for section in document.sections:
        paragraphs += section.header.paragraphs + section.footer.paragraphs
        tables += section.header.tables + section.footer.tables
    for paragraph in paragraphs:
        _update(paragraph)
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _update(paragraph)
                for subtable in cell.tables:
                    for subrow in subtable.rows:
                        for subcell in subrow.cells:
                            for subparagraph in subcell.paragraphs:
                                _update(subparagraph)
    document.save(str(path))


def update_ppt(path: pathlib.Path, data: dict[str, str]) -> None:

    def _update(container) -> None:
        if isinstance(container.text, str) and re.search("«[\w]+»", container.text):
            concatenate = False
            for index in reversed(range(len(container.runs))):
                if concatenate:
                    container.runs[index].text += container.runs[index + 1].text
                    container.runs[index + 1].text = ""
                if "»" in container.runs[index].text:
                    concatenate = True
                if "«" in container.runs[index].text:
                    concatenate = False
                    container.runs[index].text = data[container.runs[index].text]

    presentation = pptx.Presentation(str(path))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    _update(paragraph)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            _update(paragraph)
    presentation.save(str(path))


def update(directory_tpl: pathlib.Path, directory_dst: pathlib.Path, data: dict[str, str]) -> None:

    names = [path.name for path in directory_tpl.glob("*") if path.suffix in (".xlsx", ".docx", ".pptx")]

    directory_dst.mkdir()
    for name in names:
        path_src = directory_tpl / name
        path_dst = directory_dst / name
        shutil.copyfile(str(path_src), str(path_dst))

    suffix_to_function = {".xlsx": update_xls, ".docx": update_doc, ".pptx": update_ppt}

    paths = [path for path in directory_dst.glob("*") if path.suffix in (".xlsx", ".docx", ".pptx")]

    for path in paths:
        function = suffix_to_function[path.suffix]
        function(path, data)
        print(f"UPDATED | {path.name}")
